import base64
import glob
import os
import random
import re
import string
import tkinter as tk
from urllib.parse import urlparse

import openai
from icrawler import ImageDownloader
from icrawler.builtin import GoogleImageCrawler
from pptx import Presentation

unique_image_name = None


def main():
    global unique_image_name
    unique_image_name = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits) for _ in
                                range(16))

    def refresh_unique_image_name():
        global unique_image_name
        unique_image_name = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits)
                                    for _ in range(16))
        return

    class PrefixNameDownloader(ImageDownloader):

        # def get_filename(self, task, default_ext):
        #     filename = super(PrefixNameDownloader, self).get_filename(
        #         task, default_ext)
        #     refresh_unique_image_name()
        #     print(unique_image_name)
        #     return 'prefix_' + unique_image_name + filename

        def get_filename(self, task, default_ext):
            url_path = urlparse(task['file_url'])[2]
            if '.' in url_path:
                extension = url_path.split('.')[-1]
                if extension.lower() not in [
                    'jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif', 'ppm', 'pgm'
                ]:
                    extension = default_ext
            else:
                extension = default_ext
            print(unique_image_name)
            filename = base64.b64encode(url_path.encode()).decode()
            return "p_" + unique_image_name + '{}.{}'.format(filename, extension)

    def generate_ppt(topic, slide_length, api_key):
        root = Presentation("theme0.pptx")

        openai.api_key = api_key

        message = f"""Create an outline for a slideshow presentation on the topic of {topic} which is {slide_length}
        slides long. Make sure it is {slide_length} long.

        You are allowed to use the following slide types:
        Title Slide - (Title, Subtitle)
        Content Slide - (Title, Content)
        Image Slide - (Title, Content, Image)
        Thanks Slide - (Title)

        Put this tag before the Title Slide: [L_TS]
        Put this tag before the Content Slide: [L_CS]
        Put this tag before the Image Slide: [L_IS]
        Put this tag before the Thanks Slide: [L_THS]
        
        Put this tag before the Title: [TITLE]
        Put this tag after the Title: [/TITLE]
        Put this tag before the Subitle: [SUBTITLE]
        Put this tag after the Subtitle: [/SUBTITLE]
        Put this tag before the Content: [CONTENT]
        Put this tag after the Content: [/CONTENT]
        Put this tag before the Image: [IMAGE]
        Put this tag after the Image: [/IMAGE]

        Put "[SLIDEBREAK]" after each slide 

        For example:
        [L_TS]
        [TITLE]Among Us[/TITLE]

        [SLIDEBREAK]

        [L_CS] 
        [TITLE]What Is Among Us?[/TITLE]
        [CONTENT]
        1. Among Us is a popular online multiplayer game developed and published by InnerSloth.
        2. The game is set in a space-themed setting where players take on the roles of Crewmates and Impostors.
        3. The objective of Crewmates is to complete tasks and identify the Impostors among them, while the Impostors' goal is to sabotage the spaceship and eliminate the Crewmates without being caught.
        [/CONTENT]

        [SLIDEBREAK]


        Elaborate on the Content, provide as much information as possible.
        REMEMBER TO PLACE a [/CONTENT] at the end of the Content.
        Do not include any special characters (?, !, ., :, ) in the Title.
        Do not include any additional information in your response and stick to the format."""

        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {
                    "role": "user",
                    "content": message,
                },
            ],
        )

        # """ Ref for slide types:
        # 0 -> title and subtitle
        # 1 -> title and content
        # 2 -> section header
        # 3 -> two content
        # 4 -> Comparison
        # 5 -> Title only
        # 6 -> Blank
        # 7 -> Content with caption
        # 8 -> Pic with caption
        # """

        def delete_all_slides():
            for i in range(len(root.slides) - 1, -1, -1):
                r_id = root.slides._sldIdLst[i].rId
                root.part.drop_rel(r_id)
                del root.slides._sldIdLst[i]

        def create_title_slide(title, subtitle):
            layout = root.slide_layouts[0]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = subtitle

        def create_section_header_slide(title):
            layout = root.slide_layouts[2]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title

        def create_title_and_content_slide(title, content):
            layout = root.slide_layouts[1]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = content

        def create_title_and_content_and_image_slide(title, content, image_query):
            layout = root.slide_layouts[8]
            slide = root.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[2].text = content
            refresh_unique_image_name()
            google_crawler = GoogleImageCrawler(downloader_cls=PrefixNameDownloader, storage={'root_dir': os.getcwd()})
            google_crawler.crawl(keyword=image_query, max_num=1)
            dir_path = os.path.dirname(os.path.realpath(__file__))
            file_name = glob.glob(f"p_{unique_image_name}*")
            print(file_name)
            img_path = os.path.join(dir_path, file_name[0])
            slide.shapes.add_picture(img_path, slide.placeholders[1].left, slide.placeholders[1].top,
                                     slide.placeholders[1].width, slide.placeholders[1].height)

        def find_text_in_between_tags(text, start_tag, end_tag):
            start_pos = text.find(start_tag)
            end_pos = text.find(end_tag)
            result = []
            while start_pos > -1 and end_pos > -1:
                text_between_tags = text[start_pos + len(start_tag):end_pos]
                result.append(text_between_tags)
                start_pos = text.find(start_tag, end_pos + len(end_tag))
                end_pos = text.find(end_tag, start_pos)
            res1 = "".join(result)
            res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
            if len(result) > 0:
                return res2
            else:
                return ""

        def search_for_slide_type(text):
            tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
            found_text = next((s for s in tags if s in text), None)
            return found_text

        def parse_response(reply):
            list_of_slides = reply.split("[SLIDEBREAK]")
            for slide in list_of_slides:
                slide_type = search_for_slide_type(slide)
                if slide_type == "[L_TS]":
                    create_title_slide(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                       find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
                elif slide_type == "[L_CS]":
                    create_title_and_content_slide(
                        "".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                        "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                          "[/CONTENT]")))
                elif slide_type == "[L_IS]":
                    create_title_and_content_and_image_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]",
                                                                                               "[/TITLE]")),
                                                             "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                               "[/CONTENT]")),
                                                             "".join(find_text_in_between_tags(str(slide), "[IMAGE]",
                                                                                               "[/IMAGE]")))
                elif slide_type == "[L_THS]":
                    create_section_header_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

        def find_title():
            return root.slides[0].shapes.title.text

        delete_all_slides()

        print(response)

        parse_response(response.choices[0].message.content)

        root.save(fr"{find_title()}.pptx")

        print("done")

        return rf"Done! {find_title()} is ready! You can find it at {os.getcwd()}\{find_title()}.pptx"

    def button_click():
        if input2.get().isdigit():
            result_label.config(text=generate_ppt(input1.get(), input2.get(), input0.get()))

    window = tk.Tk()
    window.title("ChatGPT Generated PPTs!")
    window.configure(padx=100, pady=100)

    label0 = tk.Label(window, text="OpenAI API Key: (if this is invalid, this program will fail)")
    label0.pack()

    input0 = tk.Entry(window)
    input0.pack()

    label1 = tk.Label(window, text="Write me a PPT presentation about...")
    label1.pack()

    input1 = tk.Entry(window)
    input1.pack()

    label2 = tk.Label(window, text="Number of slides: ")
    label2.pack()

    input2 = tk.Entry(window)
    input2.pack()

    button = tk.Button(window, text="Submit", command=button_click)
    button.pack()

    result_label = tk.Label(window, text="Result")
    result_label.pack()

    window.mainloop()


if __name__ == '__main__':
    main()
