import os
import re
import json
import random
import string
from pptx import Presentation

from apis.openai_api import OpenAIClient
from apis.cohere_api import CohereAPIClient
from crawlers.icrawlercrawler import ICrawlerCrawler
from utils import get_config, get_settings

settings = get_settings()

api_clients = {
    "openai": OpenAIClient,
    "cohere": CohereAPIClient
}

image_crawlers = {
    "icrawler": ICrawlerCrawler
}


def get_generative_api_client(_generative_api_name, api_key, _generative_model_name):
    if _generative_api_name == "openai":
        return OpenAIClient(api_key, _generative_model_name)
    elif _generative_api_name == "cohere":
        return CohereAPIClient(api_key, _generative_model_name)


def generate_ppt(topic, api_name, model_name, num_slides):
    config = get_config()
    legal_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')

    save_dir = os.path.join(config["save_location"], legal_topic)
    os.makedirs(save_dir, exist_ok=True)

    ppt = Presentation("theme0.pptx")
    final_prompt = f"""Create an outline for a slideshow presentation on the topic of {topic} which is {num_slides}
        slides long. Make sure there are ONLY {num_slides} slides. This includes the title and thanks slide.

        You are allowed to use the following slide types:
        Title Slide - (Title, Subtitle)
        Content Slide - (Title, Content)
        Image Slide - (Title, Content, Image)
        Thanks Slide - (Title)
        
        Put this tag before the Title Slide: [L_TS]
        Put this tag before the Content Slide: [L_CS]
        Put this tag before the Image Slide: [L_IS]
        Put this tag before the Thanks Slide: [L_THS]
        
        Put this tag before the Title: [TITLE]
        Put this tag after the Title: [/TITLE]
        Put this tag before the Subtitle: [SUBTITLE]
        Put this tag after the Subtitle: [/SUBTITLE]
        Put this tag before the Content: [CONTENT]
        Put this tag after the Content: [/CONTENT]
        Put this tag before the Image: [IMAGE]
        Put this tag after the Image: [/IMAGE]

        Put "[SLIDEBREAK]" after each slide 

        For example:
        [L_TS]
        [TITLE]Among Us[/TITLE]

        [SLIDEBREAK]

        [L_CS] 
        [TITLE]What Is Among Us?[/TITLE]
        [CONTENT]
        1. Among Us is a popular online multiplayer game developed and published by InnerSloth.
        2. The game is set in a space-themed setting where players take on the roles of Crewmates and Impostors.
        3. The objective of Crewmates is to complete tasks and identify the Impostors among them, while the Impostors' goal is to sabotage the spaceship and eliminate the Crewmates without being caught.
        [/CONTENT]

        [SLIDEBREAK]


        Elaborate on the Content, provide as much information as possible.
        REMEMBER TO PLACE a [/CONTENT] at the end of the Content.
        Do not include any special characters (?, !, ., :, ) in the Title.
        Do not include any additional information in your response and stick to the format."""

    # """ Ref for slide types:
    # 0 -> title and subtitle
    # 1 -> title and content
    # 2 -> section header
    # 3 -> two content
    # 4 -> Comparison
    # 5 -> Title only
    # 6 -> Blank
    # 7 -> Content with caption
    # 8 -> Pic with caption
    # """

    api_key = config.get('api_key')
    api_client = get_generative_api_client(api_name, api_key, model_name)

    presentation_content = api_client.generate(final_prompt)

    ppt = Presentation("theme0.pptx")

    def delete_all_slides():
        for i in range(len(ppt.slides) - 1, -1, -1):
            r_id = ppt.slides._sldIdLst[i].rId
            ppt.part.drop_rel(r_id)
            del ppt.slides._sldIdLst[i]

    def create_title_slide(title, subtitle):
        layout = ppt.slide_layouts[0]
        slide = ppt.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def create_section_header_slide(title):
        layout = ppt.slide_layouts[2]
        slide = ppt.slides.add_slide(layout)
        slide.shapes.title.text = title

    def create_title_and_content_slide(title, content):
        layout = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    def create_title_and_content_and_image_slide(title, content, image_query):
        layout = ppt.slide_layouts[8]
        slide = ppt.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[2].text = content
        crawler = image_crawlers["icrawler"](browser="google") # TODO: Make user be able to choose crawler
        image_name = crawler.get_image(image_query, save_dir)
        img_path = os.path.join(save_dir, image_name)
        print(img_path)
        slide.shapes.add_picture(img_path, slide.placeholders[1].left, slide.placeholders[1].top,
                                 slide.placeholders[1].width, slide.placeholders[1].height)

    def find_text_in_between_tags(text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []
        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)
        res1 = "".join(result)
        res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
        if len(result) > 0:
            return res2
        else:
            return ""

    def search_for_slide_type(text):
        tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text

    def parse_response(reply):
        list_of_slides = reply.split("[SLIDEBREAK]")
        for slide in list_of_slides:
            slide_type = search_for_slide_type(slide)
            if slide_type == "[L_TS]":
                create_title_slide(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                   find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
            elif slide_type == "[L_CS]":
                create_title_and_content_slide(
                    "".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                    "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                      "[/CONTENT]")))
            elif slide_type == "[L_IS]":
                create_title_and_content_and_image_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]",
                                                                                           "[/TITLE]")),
                                                         "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                           "[/CONTENT]")),
                                                         "".join(find_text_in_between_tags(str(slide), "[IMAGE]",
                                                                                           "[/IMAGE]")))
            elif slide_type == "[L_THS]":
                create_section_header_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

    def find_title():
        return ppt.slides[0].shapes.title.text

    delete_all_slides()

    parse_response(presentation_content)

    title = "".join(str(find_title()).split(":"))

    ppt.save(os.path.join(save_dir, f"{title}.pptx"))

    return rf"Done! {title} is ready! You can find it at {os.path.realpath(save_dir)}\{title}.pptx"
